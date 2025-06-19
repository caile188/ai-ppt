#!/usr/bin/python3
# -*- coding: utf-8 -*-            
# @Author :le
# @Time : 2025/4/30 11:06

from pptx import Presentation
from pptx.util import Inches
from util.log_manager import log_manager
from PIL import Image
import os

logger = log_manager.get_logger(__name__)


def create_ppt(powerpoint_data, template_path, output_path):
    """根据解析结果生成PPTX文件"""
    prs = Presentation(template_path)
    prs.core_properties.title = powerpoint_data.title  # 设置 PowerPoint 的核心标题

    for slide_data in powerpoint_data.slides:
        slide_layout = prs.slide_layouts[slide_data.layout_id]
        new_slide = prs.slides.add_slide(slide_layout)

        # 设置幻灯片标题
        if new_slide.shapes.title:
            if slide_data.content.title:
                new_slide.shapes.title.text = slide_data.content.title
                logger.info(f"设置幻灯片标题：{slide_data.content.title}")
            elif slide_data.content.closing:
                new_slide.shapes.title.text = slide_data.content.closing
                logger.info("设置幻灯片结束页")

        # 设置文本内容
        for shape in new_slide.shapes:
            if shape.has_text_frame and not shape == new_slide.shapes.title and shape.placeholder_format.type != 18:
                tf = shape.text_frame
                tf.clear()
                first_paragraph = tf.paragraphs[0]
                if slide_data.content.presenter_name:  # 设置主讲人
                    first_paragraph.level = 0
                    format_text(first_paragraph, slide_data.content.presenter_name)
                    logger.info("设置主讲人")

                for item in slide_data.content.agenda:  # 设置目录
                    # 第一个要点覆盖初始段落，其他要点添加新段落
                    p = first_paragraph if item == slide_data.content.agenda[0] else tf.add_paragraph()
                    p.level = 0
                    format_text(p, item)  # 调用 format_text 方法来处理加粗文本
                    logger.info(f"添加目录项：{item}")

                for item in slide_data.content.bullet_points:  # 设置要点及详细阐述
                    p = first_paragraph if item == slide_data.content.bullet_points[0] else tf.add_paragraph()
                    p.level = item["level"]
                    format_text(p, item["text"])  # 调用 format_text 方法来处理加粗文本
                    logger.info(f"添加列表项：{item['text']}")

                break

        # 插入图片
        if slide_data.content.image_path:
            insert_image_centered_in_placeholder(new_slide, slide_data.content.image_path)

    prs.save(output_path)
    logger.info(f"演示文稿已保存到：{output_path}")


def format_text(paragraph, text):
    """
    格式化文本，处理加粗内容，** 包围的文本表示需要加粗。
    """
    while '**' in text:
        start = text.find('**')
        end = text.find('**', start + 2)

        if start != -1 and end != -1:
            # 添加加粗之前的普通文本
            if start > 0:
                run = paragraph.add_run()
                run.text = text[:start]

            # 添加加粗文本
            bold_run = paragraph.add_run()
            bold_run.text = text[start + 2:end]
            bold_run.font.bold = True  # 设置加粗

            # 处理剩余文本
            text = text[end + 2:]
        else:
            break

    # 添加剩余的普通文本
    if text:
        run = paragraph.add_run()
        run.text = text


def insert_image_centered_in_placeholder(new_slide, image_path):
    """
    将图片插入到 Slide 中，使其中心与 placeholder 的中心对齐。
    如果图片尺寸超过 placeholder，则进行缩小适配。
    在插入成功后删除 placeholder。
    """
    # 构建图片的绝对路径
    image_full_path = os.path.join(os.getcwd(), image_path)

    # 检查图片是否存在
    if not os.path.exists(image_full_path):
        logger.warning(f"图片路径 '{image_full_path}' 不存在，跳过此图片。")
        return

    # 打开图片并获取其大小（以像素为单位）
    with Image.open(image_full_path) as img:
        img_width_px, img_height_px = img.size

    # 遍历找到图片的 placeholder（type 18 表示图片 placeholder）
    for shape in new_slide.placeholders:
        if shape.placeholder_format.type == 18:
            placeholder_width = shape.width
            placeholder_height = shape.height
            placeholder_left = shape.left
            placeholder_top = shape.top

            # 计算 placeholder 的中心点
            placeholder_center_x = placeholder_left + placeholder_width / 2
            placeholder_center_y = placeholder_top + placeholder_height / 2

            # 图片的宽度和高度转换为 PowerPoint 的单位 (Inches)
            img_width = Inches(img_width_px / 96)  # 假设图片 DPI 为 96
            img_height = Inches(img_height_px / 96)

            # 如果图片的宽度或高度超过 placeholder，按比例缩放图片
            if img_width > placeholder_width or img_height > placeholder_height:
                scale = min(placeholder_width / img_width, placeholder_height / img_height)
                img_width *= scale
                img_height *= scale

            # 计算图片左上角位置，使其中心对准 placeholder 中心
            left = placeholder_center_x - img_width / 2
            top = placeholder_center_y - img_height / 2

            # 插入图片到指定位置并设定缩放后的大小
            new_slide.shapes.add_picture(image_full_path, left, top, width=img_width, height=img_height)
            logger.info(f"图片已插入，并以 placeholder 中心对齐，路径: {image_full_path}")

            # 移除占位符
            sp = shape._element  # 获取占位符的 XML 元素
            sp.getparent().remove(sp)  # 从父元素中删除
            logger.info("已删除图片的 placeholder")
            break